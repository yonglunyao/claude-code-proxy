import hashlib

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from ppt_extend.api import PresentationExtend


def add_aigc_mark(prs: Presentation, aigc_signature: str, request_id: str, add_visible_mark: bool = True) -> None:
    try:
        # 只在第一页添加"内容由ai生成"标识（如果 add_visible_mark 为 True）
        if add_visible_mark:
            add_ai_watermark_first_slide(prs, request_id)
        prs_ex = PresentationExtend(prs)
        all_text = extract_all_text_from_ppt(prs, request_id=request_id)
        sha256_result = calculate_sha256(all_text)
        custom_properties_part = prs_ex.custom_properties_part
        pid = custom_properties_part.next_id
        custom_properties = custom_properties_part.custom_properties
        custom_property_fmtid = 'D5CDD505-2E9C-101B-9397-08002B2CF9AE'
        custom_properties.add_property("AIGC", aigc_signature, custom_property_fmtid, pid)
    except Exception as e:
        print(f'request_id={request_id}, Update doc metadata error due to: {str(e)}')


def add_aigc_mark_to_pptx(input_path: str, output_path: str, aigc_signature: str, request_id: str, add_visible_mark: bool = True) -> None:
    prs = Presentation(input_path)
    add_aigc_mark(prs, aigc_signature, request_id, add_visible_mark)
    prs.save(output_path)


def extract_all_text_from_ppt(prs: Presentation, request_id: str = "") -> str:
    """
    从PPT文件中提取所有文字内容

    Args:
        ppt_path (str): PPT文件路径

    Returns:
        str: 拼接后的所有文字内容
    """
    try:
        # 加载PPT文件
        all_text = []

        # 遍历每一页幻灯片
        for slide_num, slide in enumerate(prs.slides, 1):
            # 遍历当前页的所有形状
            for shape in slide.shapes:
                # 跳过没有文本的形状
                if not shape.has_text_frame:
                    continue

                # 处理普通文本框
                if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text.strip():  # 跳过空文本
                                all_text.append(run.text.strip())

                # 处理占位符（标题、正文等）
                elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text.strip():  # 跳过空文本
                                all_text.append(run.text.strip())

        # 拼接所有文本，使用换行符分隔
        combined_text = '\n'.join(all_text)
        return combined_text

    except Exception as e:
        print(f"{request_id}, 读取PPT文件时出错: {e}")
        return ""


def calculate_sha256(text):
    """
    计算文本的SHA256哈希值

    Args:
        text (str): 要计算哈希的文本

    Returns:
        str: 十六进制格式的SHA256哈希值
    """
    # 编码为UTF-8（哈希计算需要字节数据）
    text_bytes = text.encode('utf-8')
    # 创建SHA256对象并计算哈希
    sha256_hash = hashlib.sha256(text_bytes).hexdigest()
    return sha256_hash


def add_ai_watermark_first_slide(prs, request_id):
    """只在PPT的第一页底部右侧位置添加'内容由ai生成'标识"""
    if len(prs.slides) == 0:
        print(f"request_id={request_id}, no slides found in presentation")
        return

    # 只处理第一页幻灯片
    slide = prs.slides[0]

    # 计算底部右侧位置
    # 宽度设为3英寸，右侧留0.2英寸边距
    left = prs.slide_width - Inches(3.2)
    # 距离底部0.2英寸
    top = prs.slide_height - Inches(0.5)

    # 添加文本框
    textbox = slide.shapes.add_textbox(left, top, Inches(3), Inches(0.3))
    tf = textbox.text_frame
    tf.text = "内容由AI生成"

    # 设置文本格式
    p = tf.paragraphs[0]
    # 右对齐
    p.alignment = PP_ALIGN.RIGHT

    # 设置字体样式
    run = p.runs[0]
    # 字体大小
    run.font.size = Pt(16)
    # 灰色字体
    run.font.color.rgb = RGBColor(128, 128, 128)
    run.font.name = "微软雅黑"

    print(f"request_id={request_id}, added AI watermark at bottom right of first slide")
